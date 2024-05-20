from fastapi import FastAPI, Request, HTTPException
from pptx import Presentation
import os, shutil

app = FastAPI()

@app.post("/generate/")
async def generate_presentation(request: Request):
    data = await request.json()
    photos_folder_path = data.get("photos_folder_path")

    if not isinstance(photos_folder_path, str):
        raise HTTPException(status_code=400, detail="Invalid 'photos_folder_path' in request data")

    if not os.path.isdir(photos_folder_path):
        raise HTTPException(status_code=400, detail=f"Invalid folder path: '{photos_folder_path}'")

    photo_files = sorted(os.listdir(photos_folder_path))
    if not photo_files:
        raise HTTPException(status_code=400, detail=f"No photo files found in folder: '{photos_folder_path}'")

    # Создаем презентацию
    prs = Presentation()
    for img_name in photo_files:
        img_path = os.path.join(photos_folder_path, img_name)
        if os.path.isfile(img_path):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = img_name
            slide.shapes.add_picture(img_path, left=0, top=0, width=prs.slide_width, height=prs.slide_height)

    presentation_path = 'C:/DPO/prezi/presentation.pptx'
    prs.save(presentation_path)

    shutil.rmtree(photos_folder_path)
    
    if os.path.isfile(presentation_path):
        return {"presentation_url": presentation_path}
    else:
        raise HTTPException(status_code=500, detail="Failed to save presentation")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8002)
